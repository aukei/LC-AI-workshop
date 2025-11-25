#!/usr/bin/env python3
import re
import sys
import subprocess
import tempfile
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

MD_PATH = Path("cognitive-thinking-ai-presentation.md")
OUT_PATH = Path("Cognitive-Thinking-in-AI.pptx")
DIAGRAM_DIR = Path("assets/diagrams")

link_pat = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")
bold_pat = re.compile(r"(\*\*|__)(.*?)\1")
ital_pat = re.compile(r"(\*|_)(.*?)\1")

class SlideData:
    def __init__(self, num, header, title, lines):
        self.num = num
        self.header = header
        self.title = title
        self.lines = lines
        self.notes = []
        self.bullets = []
        self.diagrams = []  # list[str] mermaid sources

def md_to_plain(s):
    s = link_pat.sub(lambda m: f"{m.group(1)} ({m.group(2)})", s)
    s = bold_pat.sub(lambda m: m.group(2), s)
    s = ital_pat.sub(lambda m: m.group(2), s)
    return s.strip()

def parse_markdown(md_text):
    lines = md_text.splitlines()
    slides = []
    in_mermaid = False
    cur = None
    cur_title = None
    cur_header = None
    cur_num = None
    buf = []
    mermaid_buf = []
    for raw in lines:
        line = raw.rstrip()
        # Fences
        if line.strip().startswith("```"):
            fence = line.strip()[3:]
            # Opening a mermaid block
            if not in_mermaid and fence.startswith("mermaid"):
                in_mermaid = True
                mermaid_buf = []
                continue
            # Closing a block
            if in_mermaid:
                in_mermaid = False
                mermaid_buf = []
                continue
        if in_mermaid:
            mermaid_buf.append(line)
            continue
        # New slide detection
        if line.startswith("## Slide "):
            if cur is not None:
                slides.append(SlideData(cur_num, cur_header, cur_title, buf))
            buf = []
            cur_header = line
            m = re.match(r"## Slide\s+(\d+)\s*:?\s*(.*)", line)
            cur_num = int(m.group(1)) if m else None
            cur_title = None
            cur = True
            continue
        # First H3 after slide header is slide title
        if line.startswith("### ") and cur is not None and cur_title is None:
            cur_title = line[4:].strip()
            continue
        # Slide content buffer
        if cur is not None:
            if line.startswith("---"):
                continue
            if line.strip() == "":
                buf.append("")
            else:
                buf.append(line)
    if cur is not None:
        slides.append(SlideData(cur_num, cur_header, cur_title, buf))
    # Transfer diagrams collected during parse to each slide object created
    # Note: We created SlideData instances only when we appended; diagrams were collected on 'cur' temp.
    # To keep compatibility, after the loop we need to re-parse to attach diagrams.
    # Instead, re-run lightweight pass to attach diagrams based on headers within md_text (simpler: store during parse using a temp store)
    # But we already stored diagrams in 'cur' which was True sentinel, not object. So we need to re-parse minimally:
    # Easiest approach: do a second pass to collect diagrams per "## Slide N"
    slides_by_num = {s.num: s for s in slides if s.num is not None}

    # Second pass to collect diagrams accurately
    in_mermaid = False
    current_slide_num = None
    for raw in lines:
        line = raw.rstrip()
        if line.startswith("## Slide "):
            m = re.match(r"## Slide\s+(\d+)\s*:?\s*(.*)", line)
            current_slide_num = int(m.group(1)) if m else None
            continue
        if line.strip().startswith("```"):
            fence = line.strip()[3:]
            if not in_mermaid and fence.startswith("mermaid"):
                in_mermaid = True
                mermaid_buf = []
                continue
            if in_mermaid:
                in_mermaid = False
                if current_slide_num in slides_by_num and mermaid_buf:
                    diagram_src = "\n".join(mermaid_buf).strip()
                    if diagram_src:
                        slides_by_num[current_slide_num].diagrams.append(diagram_src)
                mermaid_buf = []
                continue
        if in_mermaid:
            mermaid_buf.append(line)
            continue
    return slides

def classify_content(slide):
    bullets = []
    notes = []
    for ln in slide.lines:
        if ln.startswith("- "):
            bullets.append(md_to_plain(ln[2:]))
        elif re.match(r"\d+\.\s", ln):
            bullets.append(md_to_plain(re.sub(r"^\d+\.\s*", "", ln)))
        else:
            text = md_to_plain(ln)
            if text:
                if len(text) > 180:
                    notes.append(text)
                else:
                    bullets.append(text)
    slide.bullets = [b for b in bullets if b.strip()]
    slide.notes = notes

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    else:
        slide.placeholders[1].text = ""
    return slide

def add_content_slide(prs, title, bullets, notes=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title if title else ""
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    first = True
    for b in bullets[:20]:
        if first:
            p = body.paragraphs[0]
            first = False
        else:
            p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
    if notes:
        ns = slide.notes_slide
        ntf = ns.notes_text_frame
        ntf.text = "\n".join(notes)
    return slide

def add_diagram_slide(prs, base_title, image_path):
    # Prefer "Title Only" (index 5), fallback to "Blank" (index 6) or "Title and Content" (index 1)
    layout_idx_order = [5, 6, 1]
    slide_layout = None
    for idx in layout_idx_order:
        if idx < len(prs.slide_layouts):
            slide_layout = prs.slide_layouts[idx]
            break
    slide = prs.slides.add_slide(slide_layout)
    title_text = f"{base_title} (Diagram)" if base_title else "Diagram"
    # Set title if layout has title placeholder
    try:
        if hasattr(slide, "shapes") and slide.shapes.title:
            slide.shapes.title.text = title_text
    except Exception:
        # title placeholder may not exist (blank layout)
        pass
    # Place image
    left = Inches(1)
    top = Inches(1.2)
    max_width = prs.slide_width - Inches(2)
    # Insert with width, let height auto-scale
    slide.shapes.add_picture(str(image_path), left, top, width=max_width)
    return slide

def ensure_diagram_dir():
    DIAGRAM_DIR.mkdir(parents=True, exist_ok=True)

def render_mermaid_to_png(diagram_src, out_png: Path):
    """
    Render Mermaid diagram text to PNG using mermaid-cli via npx.
    Returns True on success, False otherwise.
    """
    # Write to temp file
    with tempfile.TemporaryDirectory() as td:
        src_path = Path(td) / "diagram.mmd"
        src_path.write_text(diagram_src, encoding="utf-8")
        cmd = [
            "npx",
            "-y",
            "@mermaid-js/mermaid-cli",
            "-i",
            str(src_path),
            "-o",
            str(out_png),
            "--backgroundColor",
            "white",
        ]
        try:
            res = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
            if res.returncode == 0 and out_png.exists():
                return True
            else:
                # Print stderr for diagnostics
                sys.stderr.write(f"mmdc failed (exit={res.returncode}): {res.stderr.decode('utf-8', errors='ignore')}\n")
                return False
        except FileNotFoundError:
            sys.stderr.write("npx not found; cannot render Mermaid via mermaid-cli.\n")
            return False

def main():
    if not MD_PATH.exists():
        print(f"Markdown not found: {MD_PATH}")
        sys.exit(1)
    md_text = MD_PATH.read_text(encoding="utf-8")
    slides = parse_markdown(md_text)
    for s in slides:
        classify_content(s)

    ensure_diagram_dir()

    prs = Presentation()
    # Title slide based on Slide 1 content
    title_slide = next((s for s in slides if s.num == 1), None)
    title = "Cognitive Thinking in AI"
    subtitle_lines = []
    if title_slide:
        if title_slide.title:
            title = title_slide.title
        for l in title_slide.lines:
            t = md_to_plain(l)
            if t:
                subtitle_lines.append(t)
    # Extract presenter/date/duration
    presenter = None
    date_line = None
    dur = None
    for t in subtitle_lines:
        if t.lower().startswith("presenter:"):
            presenter = t
        elif t.lower().startswith("date:"):
            date_line = t
        elif t.lower().startswith("duration:"):
            dur = t
    subtitle = " • ".join([x for x in [presenter, date_line, dur] if x]) or "Generated on " + datetime.now().strftime("%Y-%m-%d")
    add_title_slide(prs, title, subtitle)

    # Add other slides and render diagrams
    for s in sorted([x for x in slides if x.num and x.num != 1], key=lambda x: x.num):
        stitle = s.title or re.sub(r"^##\s*Slide\s*\d+\s*:?\s*", "", s.header or "").strip() or f"Slide {s.num}"
        # If no bullets but notes exist, promote notes to bullets for display
        if not s.bullets and s.notes:
            s.bullets = s.notes[:10]
            s.notes = []
        add_content_slide(prs, stitle, s.bullets, s.notes)

        # Render and add diagrams, if any
        if getattr(s, "diagrams", None):
            for idx, diagram_src in enumerate(s.diagrams, start=1):
                out_png = DIAGRAM_DIR / f"slide{s.num:02d}_diagram{idx:02d}.png"
                success = render_mermaid_to_png(diagram_src, out_png)
                if success:
                    add_diagram_slide(prs, stitle, out_png)
                else:
                    # Add a small fallback slide noting failure
                    fail_title = f"{stitle} (Diagram rendering failed)"
                    add_content_slide(prs, fail_title, [f"Mermaid diagram could not be rendered automatically.", "Ensure Node.js and npx are installed to enable rendering."])

    prs.save(OUT_PATH.as_posix())
    print(f"Wrote {OUT_PATH}")

if __name__ == "__main__":
    main()